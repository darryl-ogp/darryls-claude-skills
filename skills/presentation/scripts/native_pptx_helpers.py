#!/usr/bin/env python3
"""Reusable python-pptx helpers for building a NATIVE, text-editable PPTX --
every title, bullet, and table cell is a real PowerPoint shape/run, not a
screenshot. Use this path when the user says "editable", "I want to edit
this in PowerPoint", "native text", or otherwise rejects the image-flat
export.sh pipeline.

This file is a library, not a standalone script: copy the functions you
need into a per-deck build script (see the __main__ demo below for the
pattern), fill in your own content, then run it.

Hard-won gotchas from building a 9-slide DS-facing deck this way (see
../references/native-pptx-gotchas.md for the full writeup):

1. TEXTBOX height is NOT a clipping bound. It's just an initial hint --
   text silently overflows past it if the content is longer than the box.
   There is no auto-shrink and no clipping. You must estimate the real
   line count from (text length, font size, box width) and size/position
   boxes accordingly, then verify by rendering (see render_pptx_preview.sh).

2. TABLE ROW height is ALSO just a hint -- with a catch that bit us twice:
   if you pass a big total table `height` and don't set each row's height
   explicitly, PowerPoint/LibreOffice stretch every row EVENLY across that
   total, creating large dead gaps in short rows. If you pass a total
   height that's too SMALL, rows silently grow past it to fit content,
   which can push the table off the bottom of the slide.
   THE FIX: always pass `row_heights` (a list, one per row) computed from
   actual content -- never rely on the table `height` argument alone.

3. Tables have NO visible gridlines by default. Call `_set_cell_borders`
   on every cell (add_table below does this for you) or the table reads
   as loose text in a grid, not a table.

4. Dense verbatim content (many bullets across many columns) cannot be
   both large-font AND fit one slide. Don't fight this by shrinking below
   ~8pt or overflowing the slide -- pick a font in the 8-11pt range for
   that one dense table and accept it reads smaller than the rest of the
   deck. Ordinary tables (3-6 short rows) can comfortably run 13-16pt.

5. There is no way to preview a PPTX without rendering it. Always verify
   with the LibreOffice + pdftoppm pipeline (render_pptx_preview.sh) after
   every build, not just the first one -- re-verify after every edit that
   touches sizing or positioning. Never ship an unrendered guess.

Font: this uses Arial throughout (Helvetica/Arial-Bold is a safe universal
default that renders identically everywhere -- Google-Fonts-only faces
like Archivo will silently fall back to a system font on a machine that
doesn't have them installed, which is fine for body text but means don't
promise a specific custom look for a native PPTX the way you can for an
HTML/PNG export).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "Arial"


def new_slide(prs, blank_layout_index=6):
    return prs.slides.add_slide(prs.slide_layouts[blank_layout_index])


def set_bg(slide, color, slide_w, slide_h):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)  # send to back
    return rect


def add_kicker(slide, text, teal, navy, white, left=Inches(0.55), top=Inches(0.35), on_dark=False):
    """Small colored square + bold all-caps label. Make this BIG (20px+) --
    a tiny kicker is how a reader loses track of which slide is about which
    product/topic in an async deck. Don't undersize it just because it's
    "just a label"."""
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Pt(5), Pt(18), Pt(18))
    sq.fill.solid()
    sq.fill.fore_color.rgb = teal
    sq.line.fill.background()
    sq.shadow.inherit = False
    tb = slide.shapes.add_textbox(left + Pt(26), top, Inches(10.5), Pt(34))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = white if on_dark else navy
    return tb


def add_title(slide, text, navy, top=Inches(0.8), size=34, color=None, left=Inches(0.55), width=Inches(12.2)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = color or navy
    return tb


def add_bullets(slide, left, top, width, height, items, navy, color, size=16, space_after=14, italic=False):
    """items: list of strings, or (bold_lead, rest) tuples for a bold lead-in."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = 1.28
        if isinstance(item, tuple):
            bold_lead, rest = item
            r1 = p.add_run()
            r1.text = "— " + bold_lead
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.name = FONT
            r1.font.color.rgb = navy
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.name = FONT
            r2.font.color.rgb = color
            r2.font.italic = italic
        else:
            r = p.add_run()
            r.text = ("— " if item and not item.startswith(" ") else "") + item
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.color.rgb = color
            r.font.italic = italic
    return tb


def _set_cell_borders(cell, color_hex="1BAE9F", width_pt=1.0):
    """PPTX tables have no visible gridlines by default -- draw a colored
    border on all 4 edges of every cell so rows/columns read as a table.
    Border elements (lnL/lnR/lnT/lnB) must be inserted BEFORE the cell's
    fill element in the tcPr schema order, hence the explicit insert index
    instead of append."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    width_emu = int(Pt(width_pt))
    for idx, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = tcPr.makeelement(qn(tag), {"w": str(width_emu), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": color_hex})
        solidFill.append(srgb)
        ln.append(solidFill)
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
        tcPr.insert(idx, ln)


def set_cell(cell, lines, navy_soft, white, size=10.5, bold=False, color=None, fill=None,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False,
             line_spacing=1.2, space_after=5, mtb=0.07):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = white
    cell.vertical_anchor = valign
    cell.margin_left = Inches(0.09)
    cell.margin_right = Inches(0.09)
    cell.margin_top = Inches(mtb)
    cell.margin_bottom = Inches(mtb)
    tf = cell.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
        r.font.color.rgb = color or navy_soft


def add_table(slide, left, top, width, height, col_widths, rows, row_heights=None, border_color="1BAE9F"):
    """rows: list of list-of-cell-specs, each a dict of kwargs for set_cell
    (build these with your own hdr()/bodycell()/labelcell() wrappers).

    row_heights: list of Inches(...), one per row -- ALWAYS pass this.
    Without it, PowerPoint stretches rows evenly across `height`, which
    creates dead gaps whenever real content is shorter than your guess
    (this was the single biggest recurring bug building a real deck this
    way -- see gotcha #2 at the top of this file). Compute each row height
    from its actual tallest cell: roughly
        chars_per_line = (col_width_in - 2*0.09) / (font_pt * 0.52 / 72)
        lines = ceil(len(text) / chars_per_line)  # per bullet, summed per cell
        row_height_in = lines * (font_pt * line_spacing / 72 + space_after_pt / 72) + 2*mtb
    then take the max across cells in that row, and re-render to confirm --
    the estimate is a starting point, not a guarantee.
    """
    nrows = len(rows)
    ncols = len(col_widths)
    gtable = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for i, w in enumerate(col_widths):
        gtable.columns[i].width = w
    if row_heights:
        for i, h in enumerate(row_heights):
            gtable.rows[i].height = h
    for ri, row in enumerate(rows):
        for ci, spec in enumerate(row):
            cell = gtable.cell(ri, ci)
            set_cell(cell, **spec)
            _set_cell_borders(cell, color_hex=border_color)
    return gtable


def add_callout(slide, left, top, width, height, text, navy, navy_soft, light_bg, teal, bold_lead=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = light_bg
    box.line.fill.background()
    box.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = teal
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.28)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    p.line_spacing = 1.25
    if bold_lead:
        r1 = p.add_run()
        r1.text = bold_lead + " "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.name = FONT
        r1.font.color.rgb = navy
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.name = FONT
    r.font.color.rgb = navy_soft
    return box


if __name__ == "__main__":
    # Minimal self-contained demo: prove the helpers work end to end.
    # Copy this pattern into a real per-deck build script.
    TEAL = RGBColor(0x1B, 0xAE, 0x9F)
    NAVY = RGBColor(0x15, 0x2A, 0x3D)
    NAVY_SOFT = RGBColor(0x3C, 0x51, 0x64)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xF2, 0xF7, 0xF6)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = new_slide(prs)
    add_kicker(s, "Demo Slide", TEAL, NAVY, WHITE)
    add_title(s, "Native PPTX helpers work", NAVY, top=Inches(0.85), size=32)
    add_bullets(s, Inches(0.55), Inches(1.75), Inches(6), Inches(2), [
        "Every word here is a real, editable PowerPoint run.",
        "Tables below have visible borders and content-sized rows.",
    ], NAVY, NAVY_SOFT, size=16)

    def hdr(lines):
        return dict(lines=lines, size=13, bold=True, color=NAVY, fill=WHITE, navy_soft=NAVY_SOFT, white=WHITE)

    def cell(lines):
        return dict(lines=lines, size=12.5, color=NAVY_SOFT, navy_soft=NAVY_SOFT, white=WHITE)

    rows = [
        [hdr(["Column A"]), hdr(["Column B"])],
        [cell(["Row 1, short"]), cell(["Row 1, also short"])],
    ]
    add_table(s, Inches(0.55), Inches(4.2), Inches(9), Inches(1.0),
              [Inches(4.5), Inches(4.5)], rows,
              row_heights=[Inches(0.4), Inches(0.4)])

    out = "/tmp/native_pptx_demo.pptx"
    prs.save(out)
    print("wrote", out, "-- verify with render_pptx_preview.sh before trusting it")
