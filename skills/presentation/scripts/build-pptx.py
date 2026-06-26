#!/usr/bin/env python3
"""Build a 16:9 PPTX from one or more PNGs (one slide per image).

Usage:
    python3 build-pptx.py <out.pptx> <slide1.png> [<slide2.png> ...]

The PPTX is a standard PowerPoint widescreen file (13.333 x 7.5 inches).
Each PNG is embedded full-bleed on a blank-layout slide. Text is not
editable in PPT — the slide is rendered as an image — but this is the
right tradeoff for design-heavy custom HTML decks where re-implementing
the layout in native PPT shapes would lose fidelity and waste time.
"""
import sys
import os
from pptx import Presentation
from pptx.util import Inches


def main(argv):
    if len(argv) < 3:
        sys.stderr.write(__doc__)
        sys.exit(1)

    out_path = argv[1]
    pngs = argv[2:]

    for p in pngs:
        if not os.path.isfile(p):
            sys.stderr.write(f"missing PNG: {p}\n")
            sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(out_path)
    size_kb = os.path.getsize(out_path) / 1024
    print(f"wrote {out_path} ({size_kb:.0f}KB, {len(pngs)} slide{'s' if len(pngs) != 1 else ''})")


if __name__ == "__main__":
    main(sys.argv)
